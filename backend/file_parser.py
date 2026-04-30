import json
import os
import shutil
import subprocess
from typing import List, Optional

from langchain_core.documents import Document
from langchain_community.document_loaders import TextLoader


def _extract_ppt_text(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    text_content = []
    for slide in prs.slides:
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_parts.append(shape.text)
        if slide_parts:
            text_content.append("\n".join(slide_parts))
    return "\n\n".join(text_content)


def _extract_docx_text(file_path: str) -> Optional[str]:
    try:
        from docx import Document as DocxDocument

        doc = DocxDocument(file_path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    parts.append(" | ".join(row_text))
        return "\n".join(parts)
    except Exception as e:
        print(f"python-docx 解析失败 {file_path}: {str(e)}")
        return None


def _extract_doc_text_with_libreoffice(file_path: str) -> Optional[str]:
    try:
        from docx2txt import process

        text = process(file_path)
        return text.strip() if text else None
    except Exception as e:
        print(f"docx2txt 解析失败 {file_path}: {str(e)}")
        return None


def _extract_doc_text_with_win32com(file_path: str) -> Optional[str]:
    try:
        import win32com.client  # type: ignore

        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(os.path.abspath(file_path))
        text = doc.Content.Text
        doc.Close(False)
        word.Quit()
        return text
    except Exception as e:
        print(f"win32com 解析 DOC 失败 {file_path}: {str(e)}")
        return None


def _extract_ppt_text_with_win32com(file_path: str) -> Optional[str]:
    try:
        import win32com.client  # type: ignore

        ppt = win32com.client.Dispatch("PowerPoint.Application")
        presentation = ppt.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
        parts = []
        for slide in presentation.Slides:
            slide_parts = []
            for shape in slide.Shapes:
                try:
                    if shape.HasTextFrame and shape.TextFrame.HasText:
                        text = shape.TextFrame.TextRange.Text
                        if text.strip():
                            slide_parts.append(text)
                except Exception:
                    continue
            if slide_parts:
                parts.append("\n".join(slide_parts))
        presentation.Close()
        ppt.Quit()
        return "\n\n".join(parts)
    except Exception as e:
        print(f"win32com 解析 PPT 失败 {file_path}: {str(e)}")
        return None


def _convert_with_libreoffice(file_path: str, target_dir: str, convert_to: str) -> Optional[str]:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", convert_to, "--outdir", target_dir, file_path],
            check=True,
            capture_output=True,
            text=True,
        )
    except Exception:
        return None
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    converted_path = os.path.join(target_dir, f"{base_name}.{convert_to}")
    return converted_path if os.path.exists(converted_path) else None


def _write_text_sidecar(file_path: str, text: str) -> Optional[str]:
    if text is None:
        return None
    txt_path = os.path.splitext(file_path)[0] + ".txt"
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(text)
    return txt_path


def convert_to_text_if_needed(file_path: str) -> Optional[str]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in {".doc", ".docx", ".ppt", ".pptx"}:
        return None
    try:
        if ext == ".docx":
            text = _extract_docx_text(file_path)
            if text:
                return _write_text_sidecar(file_path, text)
            text = _extract_doc_text_with_libreoffice(file_path)
            if text:
                return _write_text_sidecar(file_path, text)
            return None
        if ext == ".pptx":
            text = _extract_ppt_text(file_path)
            if text:
                return _write_text_sidecar(file_path, text)
            return None
        if ext == ".doc":
            text = _extract_docx_text(file_path)
            if text:
                return _write_text_sidecar(file_path, text)
            text = _extract_doc_text_with_win32com(file_path)
            if text:
                return _write_text_sidecar(file_path, text)
            text = _extract_doc_text_with_libreoffice(file_path)
            if text:
                return _write_text_sidecar(file_path, text)
            converted = _convert_with_libreoffice(file_path, os.path.dirname(file_path), "docx")
            if converted:
                return convert_to_text_if_needed(converted)
            return None
        if ext == ".ppt":
            text = _extract_ppt_text_with_win32com(file_path)
            if text:
                return _write_text_sidecar(file_path, text)
            converted = _convert_with_libreoffice(file_path, os.path.dirname(file_path), "pptx")
            if converted:
                return convert_to_text_if_needed(converted)
            return None
    except Exception as e:
        print(f"自动转换文件失败 {file_path}: {str(e)}")
    return None


def parse_file(file_path: str) -> List[Document]:
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".pdf":
            try:
                import pdfplumber
                try:
                    from rapidocr_onnxruntime import RapidOCR
                    import fitz  # PyMuPDF
                    ocr_engine = RapidOCR()
                    print("OCR 引擎初始化成功")
                except ImportError:
                    ocr_engine = None
                    print("未安装 RapidOCR 或 fitz，OCR 功能将不可用")

                docs = []
                with pdfplumber.open(file_path) as pdf:
                    # 如果需要 OCR，我们需要用 fitz 打开同一份文件
                    pdf_doc = fitz.open(file_path) if ocr_engine else None
                    
                    for i, page in enumerate(pdf.pages):
                        text = page.extract_text() or ""
                        clean_text = text.replace('\xa0', ' ').strip()
                        
                        # 如果文字太少 (例如少于 50 字)，且 OCR 引擎可用，则尝试 OCR
                        if len(clean_text) < 50 and ocr_engine and pdf_doc:
                            try:
                                print(f"第 {i+1} 页文字过少 ({len(clean_text)}字)，正在尝试 OCR 识别...")
                                pix = pdf_doc[i].get_pixmap(matrix=fitz.Matrix(2, 2))  # 2倍缩放提高识别率
                                img_bytes = pix.tobytes("png")
                                result, _ = ocr_engine(img_bytes)
                                if result:
                                    ocr_text = "\n".join([line[1] for line in result])
                                    if len(ocr_text) > len(clean_text):
                                        clean_text = ocr_text
                                        print(f"第 {i+1} 页 OCR 识别成功，获取到 {len(clean_text)} 字")
                            except Exception as ocr_err:
                                print(f"第 {i+1} 页 OCR 识别出错: {str(ocr_err)}")

                        if clean_text:
                            docs.append(Document(
                                page_content=clean_text, 
                                metadata={"source": file_path, "page": i + 1}
                            ))
                    
                    if pdf_doc:
                        pdf_doc.close()

                if docs:
                    print(f"解析成功 [PDF-OCR增强]: {os.path.basename(file_path)} -> {len(docs)} 段")
                    return docs
                return []
            except Exception as pdf_err:
                print(f"PDF 深度解析失败: {file_path}, 原因: {str(pdf_err)}")
                # 最后的兜底：尝试最基础的 pypdf
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(file_path)
                    docs = []
                    for page_index, page in enumerate(reader.pages):
                        text = page.extract_text() or ""
                        if text.strip():
                            docs.append(Document(page_content=text.strip(), metadata={"source": file_path, "page": page_index + 1}))
                    return docs
                except:
                    return []
        if ext == ".docx":
            text = _extract_docx_text(file_path)
            if text:
                docs = [Document(page_content=text, metadata={"source": file_path})]
                print(f"解析成功 [DOCX]: {os.path.basename(file_path)} -> {len(docs)} 段")
                return docs
            converted_path = convert_to_text_if_needed(file_path)
            if converted_path and os.path.exists(converted_path):
                return parse_file(converted_path)
            return []
        if ext == ".pptx":
            text = _extract_ppt_text(file_path)
            docs = [Document(page_content=text, metadata={"source": file_path})]
            print(f"解析成功 [PPTX]: {os.path.basename(file_path)} -> {len(docs)} 段")
            return docs
        if ext in [".md", ".markdown", ".txt"]:
            loader = TextLoader(file_path, encoding="utf-8")
            docs = loader.load()
            print(f"解析成功 [TEXT]: {os.path.basename(file_path)} -> {len(docs)} 段")
            return docs
        if ext == ".json":
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            documents = []
            for item in data:
                content = f"问题: {item.get('question', '')}\n回答: {item.get('answer', '')}"
                documents.append(Document(page_content=content, metadata=item.get("metadata", {})))
            print(f"解析成功 [JSON]: {os.path.basename(file_path)} -> {len(documents)} 段")
            return documents
        if ext in {".doc", ".ppt"}:
            converted_path = convert_to_text_if_needed(file_path)
            if converted_path and os.path.exists(converted_path):
                print(f"自动转换成功 [{ext.upper()}]: {os.path.basename(file_path)} -> {os.path.basename(converted_path)}")
                return parse_file(converted_path)
            print(f"警告: 暂不支持直接解析格式 {ext}, 文件: {file_path}")
            return []
        print(f"警告: 暂不支持解析格式 {ext}, 文件: {file_path}")
        return []
    except Exception as e:
        print(f"解析文件 {file_path} 出错: {str(e)}")
        return []


def load_all_knowledge(knowledge_dir: str, dataset_json_path: str) -> List[Document]:
    all_documents = []
    if os.path.exists(dataset_json_path):
        all_documents.extend(parse_file(dataset_json_path))
    if os.path.exists(knowledge_dir):
        for root, dirs, files in os.walk(knowledge_dir):
            for file in files:
                file_path = os.path.join(root, file)
                all_documents.extend(parse_file(file_path))
    return all_documents
